import os
from pathlib import Path
from pptx import Presentation
from .papago_api import translate_ppt


def translate_papago(filename):
    BASE_DIR = Path(__file__).resolve().parent.parent
    path = os.path.join(BASE_DIR, 'media/')
    prs = Presentation(path + filename)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                translated = translate_ppt(paragraph.text)
                paragraph.clear()
                
                run = paragraph.add_run()
                run.text = translated

    prs.save(path + "translated_" + filename)