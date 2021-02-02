import os
import ssl
import json
import urllib
import urllib.request
import mimetypes

from pathlib import Path
from pptx import Presentation

from django.http import HttpResponse, Http404
from django.shortcuts import render, redirect

from .forms import UploadForm
from .models import FileUpload
from config.settings import NMT_ID, NMT_PW, MEDIA_ROOT


def file_upload(request):
    """ 번역할 파일 업로드 """
    form = UploadForm()
    if request.method == 'POST':
        form = UploadForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            name = request.FILES['up_file'].name
            translate_papago(name)
                
            return render(request, 'translator/download_form.html', {'name':name})

    return render(request, 'translator/upload_form.html', {'form':form})


def file_download(request, filename):
    """ 번역된 파일 다운로드 """
    file_name = 'translated_' + filename
    file_url = MEDIA_ROOT +'/' + file_name
    
    if request.method == 'GET':
        if os.path.exists(file_url):
            with open(file_url, 'rb') as fh:
                quote_file_url = urllib.parse.quote(file_name.encode('utf-8'))
                response = HttpResponse(fh.read(), content_type=mimetypes.guess_type(file_url)[0])
                response['Content-Disposition'] = 'attachment;filename*=UTF-8\'\'%s' % quote_file_url
                return response
            raise Http404

    return render(request, 'translator/download_form.html')


def translate_papago(filename):
    """ PPT 번역 """
    file_url = MEDIA_ROOT +'/'
    prs = Presentation(file_url + filename)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                content = paragraph.text
                trimmed = ' '.join(content.split())
                if not trimmed:
                    continue
                
                translated = papago(trimmed)
                paragraph.clear()
                
                run = paragraph.add_run()
                run.text = translated

    prs.save(file_url + "translated_" + filename)


def papago(txt):
    """ 파파고 번역 API 호출 """
    nmt_id = NMT_ID
    nmt_pw = NMT_PW

    encText = urllib.parse.quote(txt)
    data = "source=ko&target=en&text=" + encText
    url = "https://openapi.naver.com/v1/papago/n2mt"
    
    request = urllib.request.Request(url)
    request.add_header("X-Naver-Client-Id", nmt_id)
    request.add_header("X-Naver-Client-Secret", nmt_pw)
    res_ssl = ssl._create_unverified_context()
    response = urllib.request.urlopen(request, data=data.encode("utf-8"), context=res_ssl)
    rescode = response.getcode()

    if(rescode==200):
        response_body = response.read()
        res = json.loads(response_body.decode('utf-8'))
        return(res['message']['result']['translatedText'])

    else:
        return("Error Code:" + rescode) 